"""Generate meeting slides for XelToFab progress meeting with Prof. Mark Fuge."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Color palette ---
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BG = RGBColor(0x16, 0x16, 0x25)
ACCENT = RGBColor(0x4E, 0x9A, 0xF5)      # blue
GREEN = RGBColor(0x4E, 0xC9, 0xB0)       # teal-green
YELLOW = RGBColor(0xDC, 0xDC, 0x8B)      # muted yellow
GRAY = RGBColor(0x88, 0x88, 0x99)        # muted gray
RED_ACCENT = RGBColor(0xE0, 0x6C, 0x75)  # soft red
SUBTITLE_CLR = RGBColor(0xAA, 0xBB, 0xCC)
BODY_CLR = RGBColor(0xCC, 0xCC, 0xDD)
LIGHT_BG = RGBColor(0x1E, 0x1E, 0x30)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, left=0.8, top=0.4, width=8.4, height=1.0,
              font_size=32, color=WHITE, bold=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return txBox


def add_subtitle(slide, text, left=0.8, top=1.3, width=8.4, height=0.6,
                 font_size=18, color=SUBTITLE_CLR):
    return add_title(slide, text, left, top, width, height, font_size, color, bold=False)


def add_body_text(slide, lines, left=0.8, top=2.0, width=8.4, height=4.5,
                  font_size=16, color=BODY_CLR, line_spacing=1.4):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle bold markers **text**
        parts = line.split("**")
        for j, part in enumerate(parts):
            if not part:
                continue
            run = p.add_run()
            run.text = part
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = (j % 2 == 1)

        p.space_after = Pt(font_size * (line_spacing - 1) + 4)

        # Indent bullet lines
        if line.startswith("  "):
            p.level = 1
    return txBox


def add_accent_bar(slide, top=1.15):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(top), Inches(2.0), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def add_flow_arrow(slide, labels, top=2.0, colors=None):
    """Add a horizontal flow diagram with arrow-connected boxes."""
    n = len(labels)
    box_w = 1.4
    gap = 0.35
    total = n * box_w + (n - 1) * gap
    start_x = (10 - total) / 2

    for i, label in enumerate(labels):
        x = start_x + i * (box_w + gap)
        clr = colors[i] if colors else ACCENT

        # Box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top), Inches(box_w), Inches(0.7)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = clr
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE
        run.font.bold = True

        # Arrow
        if i < n - 1:
            ax = x + box_w
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(ax + 0.02), Inches(top + 0.22), Inches(gap - 0.04), Inches(0.26)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()


def add_table(slide, headers, rows, left=0.8, top=2.2, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [8.4 / n_cols] * n_cols

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(sum(col_widths)), Inches(0.35 * n_rows)
    )
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    def style_cell(cell, text, is_header=False):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(12)
        run.font.bold = is_header
        run.font.color.rgb = WHITE if is_header else BODY_CLR
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if is_header else LIGHT_BG
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for j, h in enumerate(headers):
        style_cell(table.cell(0, j), h, is_header=True)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            style_cell(table.cell(i + 1, j), val)

    return table_shape


# ============================================================
# Build presentation
# ============================================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)  # 16:9
blank_layout = prs.slide_layouts[6]  # blank

# ------ Slide 1: Title ------
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title(slide, "XelToFab", top=1.2, font_size=44, color=ACCENT)
add_subtitle(slide, "From Density Fields to Fabrication-Ready Geometry",
             top=2.0, font_size=22, color=WHITE)
add_body_text(slide, [
    "M4X \u2013 IDEAL Lab Collaboration",
    "Prof. Mark Fuge",
], left=0.8, top=3.0, font_size=16, color=SUBTITLE_CLR)
add_accent_bar(slide, top=1.85)

# ------ Slide 2: The Problem ------
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title(slide, "The Post-Processing Gap")
add_accent_bar(slide)
add_body_text(slide, [
    "\u2022 Every ML-TO method (TopoDiff, NITO, NTopo) outputs **raw density fields**",
    "\u2022 EngiBench benchmarks store designs as numpy arrays \u2014 no mesh extraction",
    "\u2022 Researchers manually threshold + export in ABAQUS or ParaView",
    "\u2022 **No automated density \u2192 mesh conversion** in any existing benchmark",
    "",
], top=1.6)
add_flow_arrow(slide, [
    "TO Solver", "Density\nField", "???", "Usable\nGeometry"
], top=4.2, colors=[GREEN, ACCENT, RED_ACCENT, GRAY])

# ------ Slide 3: Pipeline Architecture ------
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title(slide, "Pipeline Architecture")
add_accent_bar(slide)
add_flow_arrow(slide, [
    "Preprocess", "Extract", "Post-process", "Export"
], top=1.7, colors=[GREEN, GREEN, YELLOW, GREEN])
add_body_text(slide, [
    "\u2022 **Preprocess**: Heaviside threshold, Gaussian smooth, morphological cleanup",
    "\u2022 **Extract**: Marching cubes (3D) / marching squares (2D)",
    "\u2022 **Post-process**: Taubin smoothing, decimation, remeshing, quality metrics",
    "\u2022 **Export**: STL / OBJ / PLY for fabrication, VTK for FEA",
    "",
    "  Green = implemented  \u2502  Yellow = next tier  \u2502  Gray = future",
], top=2.7, font_size=15)

# ------ Slide 4: What's Working Now ------
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title(slide, "What\u2019s Working Now (MVP)")
add_accent_bar(slide)
add_body_text(slide, [
    "\u2022 **Preprocessing**: threshold + smooth + morphology + connected components",
    "\u2022 **Extraction**: scikit-image marching cubes/squares",
    "\u2022 **Smoothing**: Taubin smoothing, validated volume preservation (>90%)",
    "\u2022 **I/O**: 11 input formats (NumPy, MATLAB, CSV, VTK, HDF5, XDMF)",
    "\u2022 **CLI**: xtf process density.npy -o result.stl --threshold 0.4",
    "\u2022 **Tested**: 816 LOC production, 979 LOC tests, 73 passing tests",
], top=1.6)

# ------ Slide 5: Demo / Visual Results ------
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title(slide, "Demo / Visual Results")
add_accent_bar(slide)
add_body_text(slide, [
    "Pipeline in action:",
    "",
    "\u2022 Raw density field \u2192 preprocessed binary \u2192 extracted mesh",
    "\u2022 2D contour + 3D mesh examples",
    "\u2022 Volume fraction tracking through stages",
    "",
    "[ Screenshots / live marimo demo ]",
], top=1.6)
# Placeholder boxes for screenshots
for i, label in enumerate(["Raw Density", "Binary", "Mesh"]):
    x = 1.5 + i * 2.6
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(3.4), Inches(2.2), Inches(1.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = GRAY
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = f"[ {label} ]"
    run.font.size = Pt(12)
    run.font.color.rgb = GRAY

# ------ Slide 6: TO_fixtures ------
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title(slide, "TO_fixtures: State & Opportunities")
add_accent_bar(slide)
add_body_text(slide, [
    "\u2022 Student-built SIMP solver, **Python 2.7**, locked to ABAQUS",
    "\u2022 Used by a single PhD \u2014 STL_Export.py is an **empty skeleton**",
    "\u2022 No mesh output, no standalone visualization, no quality metrics",
    "",
    "**XelToFab sets a higher bar:**",
    "  \u2022 Format-agnostic (any solver, not just ABAQUS)",
    "  \u2022 Batch CLI automation vs. manual GUI interaction",
    "  \u2022 Validated & tested pipeline vs. ad-hoc scripts",
    "",
    "\u2022 Integration: 1 line (np.save) \u2192 xtf process",
], top=1.6, font_size=15)

# ------ Slide 7: Tier 2 Next Steps ------
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title(slide, "Next: Quality-Enhanced Pipeline (Tier 2)")
add_accent_bar(slide)
add_table(slide,
    ["Feature", "Method", "Goal"],
    [
        ["QEM decimation", "quadric edge collapse (pymeshlab)", "Reduce triangle count"],
        ["Isotropic remeshing", "split/collapse/flip/smooth", "Uniform triangle quality"],
        ["Feature smoothing", "bilateral mesh filtering", "Preserve sharp edges"],
        ["Watertight repair", "hole filling + manifold fixing", "Printable meshes"],
        ["Quality metrics", "aspect ratio, Jacobian, min angle", "FEA-readiness score"],
    ],
    col_widths=[2.2, 3.4, 2.8],
    top=1.7,
)
add_body_text(slide, [
    "**Target**: meshes directly usable for re-analysis or 3D printing",
], top=4.5, font_size=15)

# ------ Slide 8: Beyond Meshes ------
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title(slide, "Beyond Meshes: When Is CAD Needed?")
add_accent_bar(slide)
add_body_text(slide, [
    "**For many workflows, STL/mesh is sufficient:**",
    "  \u2022 3D printing, CNC from mesh slicing, FEA re-analysis",
    "",
    "**CAD (NURBS/B-Rep/STEP) matters when:**",
    "  \u2022 Integration with existing CAD assemblies is required",
    "  \u2022 Parametric editing of the optimized shape is needed",
    "  \u2022 Downstream processes mandate STEP/IGES (e.g., injection molding)",
    "",
    "\u2022 **Need a concrete use case** to justify the research investment",
    "\u2022 Framing: \u201cWhen does the mesh-to-CAD gap actually block practitioners?\u201d",
], top=1.6, font_size=15)

# ------ Slide 9: Future Directions ------
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title(slide, "Possible Future Directions")
add_accent_bar(slide)
add_body_text(slide, [
    "\u2022 **Differentiable extraction**: FlexiCubes / Kaolin \u2014 mesh quality",
    "  as a training objective, end-to-end gradient flow",
    "",
    "\u2022 **Neural implicit representations**: NITO / NTopo \u2014 smoother fields,",
    "  resolution-independent, better meshes by construction",
    "",
    "\u2022 **Learned mesh-to-CAD**: Point2CAD-style patch decomposition from data",
    "",
    "\u2022 **EngiBench ecosystem**: XelToFab as independent library;",
    "  separate demo project to showcase integration",
], top=1.6, font_size=15)

# ------ Slide 10: Discussion ------
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_title(slide, "Discussion")
add_accent_bar(slide)
add_body_text(slide, [
    "**Questions for discussion:**",
    "",
    "\u2022 Should differentiable extraction (FlexiCubes) be a research direction,",
    "  or is the classical pipeline sufficient for our needs?",
    "",
    "\u2022 What TO problem classes should we prioritize for validation \u2014",
    "  compliance beams, multi-constraint mechanisms, 3D lattices?",
    "",
    "\u2022 Are there downstream users / collaborators at IDEAL Lab who would",
    "  exercise the pipeline on real problems and give feedback?",
    "",
    "\u2022 What\u2019s the right publication venue \u2014 tools/software paper (JOSS,",
    "  SoftwareX) or methods contribution tied to a TO application?",
], top=1.6, font_size=15)

# ------ Save ------
output_path = "docs/2026-03-03-meeting-slides.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
